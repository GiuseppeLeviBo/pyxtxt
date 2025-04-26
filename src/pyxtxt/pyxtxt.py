import fitz  # PyMuPDF
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
from odf.opendocument import load
from odf.text import P
import openpyxl

def xtxt_pdf(file_buffer):
    try:
        raw_data = file_buffer.read()
        if not raw_data:
            print("⚠️  PDF vuoto o non letto correttamente.")
            return None

        doc = fitz.open(stream=raw_data, filetype="pdf")
        return "\n".join(page.get_text() for page in doc)

    except fitz.EmptyFileError:
        print("⚠️  Errore: il PDF è vuoto o non leggibile.")
        return None
    except Exception as e:
        print(f"⚠️  Errore durante l'estrazione del PDF: {e}")
        return None



def xtxt_docx(file_buffer) -> str:
    try:
        # Copia del buffer per poterlo riutilizzare
        file_buffer.seek(0)
        data = file_buffer.read()
        buffer_copy = io.BytesIO(data)

        if not zipfile.is_zipfile(buffer_copy):
            print("⚠️  DOCX non valido (non è un file zip)")
            return ""

        buffer_copy.seek(0)
        doc = Document(buffer_copy)

        text = "\n".join(p.text for p in doc.paragraphs)
        return text

    except Exception as e:
        print(f"⚠️  Errore durante l'estrazione DOCX : {e}")
        return ""

import zipfile


def xtxt_pptx(file_buffer) -> str:
    try:
        # Convertiamo il file_buffer (che è già un BytesIO o simile) in modo da poterlo riusare
        file_buffer.seek(0)
        data = file_buffer.read()
        buffer_copy = io.BytesIO(data)

        if not zipfile.is_zipfile(buffer_copy):
            print("⚠️  PPTX non valido (non è un file zip)")
            return ""

        # Se è un file zip valido, possiamo ripassare i dati a Presentation
        buffer_copy.seek(0)
        prs = Presentation(buffer_copy)

        text = "\n".join(
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
        return text

    except Exception as e:
        print(f"⚠️  Errore durante l'estrazione PPTX : {e}")
        return ""


from openpyxl import load_workbook
from openpyxl.worksheet.worksheet import Worksheet
import openpyxl

def xtxt_xlsx(file_buffer, max_rows_per_sheet: int = 200) -> str:
    try:
        file_buffer.seek(0)
        data = file_buffer.read()
        buffer_copy = io.BytesIO(data)

        if not zipfile.is_zipfile(buffer_copy):
            print("⚠️  XLSX non valido (non è un archivio zip)")
            return ""

        buffer_copy.seek(0)
        wb = openpyxl.load_workbook(buffer_copy, data_only=True, read_only=True)
    except Exception as e:
        print(f"⚠️ Errore durante la lettura XLSX: {e}")
        return ""

    testo = []
    for sheet in wb.worksheets:
        if sheet.sheet_state != 'visible':
            continue
        testo.append(f"# {sheet.title}")
        count = 0
        for row in sheet.iter_rows(values_only=True):
            if max_rows_per_sheet != -1 and count >= max_rows_per_sheet:
                break
            valori = [str(cell).strip() if cell is not None else "" for cell in row]
            if any(valori):
                testo.append(" | ".join(valori))
                count += 1

    return "\n".join(testo)

def xtxt_txt(file_buffer):
    return file_buffer.read().decode("utf-8", errors="ignore")


def xtxt_odt(file_buffer):
    odt_doc = load(file_buffer)
    paragraphs = odt_doc.getElementsByType(P)
    
    testo = []
    for p in paragraphs:
        contenuto = []
        for n in p.childNodes:
            if n.nodeType == 3:  # TEXT_NODE
                contenuto.append(n.data)
        if contenuto:
            testo.append("".join(contenuto))
    
    return "\n".join(testo)


def xtxt_html(file_buffer):
    soup = BeautifulSoup(file_buffer.read(), "html.parser")
    return soup.get_text(separator="\n")





from lxml import etree

def xtxt_xml(file_buffer) -> str:
    try:
        file_buffer.seek(0)
        parser = etree.XMLParser(recover=True)
        tree = etree.parse(file_buffer, parser)
        root = tree.getroot()

        # Estrai il testo ricorsivamente da tutti i nodi
        def get_text_recursively(elem):
            texts = []
            if elem.text:
                texts.append(elem.text.strip())
            for child in elem:
                texts.append(get_text_recursively(child))
                if child.tail:
                    texts.append(child.tail.strip())
            return " ".join(filter(None, texts))

        testo = get_text_recursively(root)
        return testo.strip()

    except Exception as e:
        print(f"⚠️ Errore durante l'estrazione XML: {e}")
        return ""
import xlrd

def xtxt_xls(file_buffer, max_rows_per_sheet: int = 100) -> str:
    try:
        file_buffer.seek(0)
        workbook = xlrd.open_workbook(file_contents=file_buffer.read())
        testo = []

        for sheet in workbook.sheets():
            testo.append(f"# {sheet.name}")
            for row_idx in range(min(sheet.nrows, max_rows_per_sheet)):
                row = sheet.row(row_idx)
                valori = [str(cell.value).strip() for cell in row if str(cell.value).strip()]
                if valori:
                    testo.append(" | ".join(valori))

        return "\n".join(testo)

    except Exception as e:
        print(f"⚠️ Errore durante l'estrazione XLS: {e}")
        return ""

# from pyth.plugins.plaintext.reader import PlaintextReader
# from pyth.plugins.plaintext import PlaintextFile

# def xtxt_rtf(file_buffer):
    # try:
        # rtf = PlaintextReader.read(file_buffer)
        # return rtf
    # except Exception as e:
        # print(f"⚠️ Errore durante l'estrazione RTF: {e}")
        # return ""

from lxml import etree

def xtxt_svg(file_buffer):
    try:
        tree = etree.parse(file_buffer)
        root = tree.getroot()

        # Estrai tutto il testo dai tag <text>
        texts = [element.text for element in root.findall('.//{http://www.w3.org/2000/svg}text')]
        return "\n".join(texts)
    except Exception as e:
        print(f"⚠️ Errore durante l'estrazione SVG: {e}")
        return ""

from functools import singledispatch
import io
import magic

@singledispatch
def xtxt(file_input):
    raise NotImplementedError(f"Tipo non supportato: {type(file_input)}")

# Caso 1: file path (str)
@xtxt.register
def _(file_input: str):
    try:
        with open(file_input, "rb") as f:
            data = f.read()
        buffer = io.BytesIO(data)
        buffer.name=file_input
        buffer.mimeType=magic.Magic(mime=True).from_file(file_input)
        return xtxt(buffer)
    except Exception as e:
        print(f"⚠️ Errore apertura file '{file_input}': {e}")
        return None

# Caso 2: buffer (BytesIO)
@xtxt.register
def _(file_input: io.BytesIO):
    try:


        # Mappa dei MIME Type gestiti dagli estrattori
        estrattori = {
            "application/pdf": xtxt_pdf,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": xtxt_docx,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": xtxt_pptx,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": xtxt_xlsx,
            "application/vnd.ms-excel": xtxt_xls,
            "text/plain": xtxt_txt,
            "application/vnd.oasis.opendocument.text": xtxt_odt,
            "text/html": xtxt_html,
#            "text/rtf": xtxt_rtf,
            "application/xml": xtxt_xml,
            "text/xml": xtxt_xml,
        }
        if hasattr(file_input,'mimeType'):
            mime_type=file_input.mimeType
        else:
            mime_type = magic.Magic(mime=True).from_buffer(file_input.read(2048))
            file_input.name='IO_buffer'
            file_input.seek(0)
        print(mime_type)
        if mime_type.startswith("text/"):
            if (mime_type != "text/html") and (mime_type != "text/xml") and (mime_type != "text/plain"):
               print(f"📄 File di tipo testuale riconosciuto: {mime_type}, trattato come text/plain")
               mime_type = "text/plain"
        if mime_type not in estrattori:
            print(f"⚠️ MIME type non supportato dopo esportazione: {mime_type} ({file_input.name}) salto.")
            return None


        # Estrai il testo
        testo = estrattori[mime_type](file_input)

        # Ottieni il percorso completo del file
        print(f"✅ Testo estratto da {file_input}")
        return f"{testo}"
    except  Exception as e:
        print(f"❌ Errore durante la lettura: {e}")
        return None
